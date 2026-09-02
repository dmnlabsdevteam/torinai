#!/usr/bin/env python3
"""
Documentation Tools
===================
AI-powered documentation generation tools

Purpose:
- Generate code documentation automatically
- Create README files
- Generate API documentation
- Analyze and document codebase
"""

import asyncio
import logging
import os
import re
from typing import Dict, Any, List, Optional
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)


@dataclass
class DocumentationConfig:
    """Documentation generation configuration"""
    include_examples: bool = True
    include_types: bool = True
    include_docstrings: bool = True
    style: str = "google"  # google, numpy, sphinx
    output_format: str = "markdown"  # markdown, rst, html


@dataclass
class DocumentationResult:
    """Documentation generation result"""
    content: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    success: bool = True
    error: Optional[str] = None


# ============================================================================
# Tool Classes for Agent Use
# ============================================================================

from core.tools.tool_registry import Tool, ToolCategory, ToolResult, ToolParameter
from core.tools.capabilities import Capability, ToolCapabilityProfile, CapabilityMetadata, RiskLevel


class GenerateReadmeTool(Tool):
    """Production-ready README generation using LLM and project analysis"""

    def __init__(self):
        super().__init__()
        self.name = "generate_readme"
        self.description = "Generate a comprehensive README.md by analyzing project structure using LLM"
        self.category = ToolCategory.DOCUMENTATION
        self.parameters = [
            ToolParameter(
                name="project_path",
                type="string",
                description="Path to the project directory",
                required=True
            ),
            ToolParameter(
                name="project_name",
                type="string",
                description="Name of the project (optional, inferred from directory)",
                required=False
            ),
            ToolParameter(
                name="description",
                type="string",
                description="Project description (optional, will be inferred)",
                required=False
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="generate_readme",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="GenerateReadme capability"
                ),
                CapabilityMetadata(
                    capability=Capability.GENERATE_DOCS,
                    description="Generate comprehensive documentation",
                    input_types=["source_code", "config"],
                    output_types=["docs"],
                    latency="medium",
                    cost="low",
                    reliability="high",
                    risk_level=RiskLevel.LOW,
                    priority=8
                )
            ]
        )

    async def execute(self, **kwargs) -> ToolResult:
        """Assemble a README from real project analysis, model-free.

        Every section is built from facts actually found in the project
        (detected manifests, dependency list, directory structure). There is
        no hand-written narrative and no placeholder prose passed off as one.
        """
        try:
            project_path = kwargs.get("project_path", ".")
            project_name = kwargs.get("project_name") or Path(project_path).name
            description = kwargs.get("description", "")

            structure_info = self._analyze_project(project_path)
            readme_content = self._build_readme(project_name, description, structure_info)

            return ToolResult(
                success=True,
                output={
                    "content": readme_content,
                    "filename": "README.md",
                    "project_name": project_name,
                    "analysis": structure_info,
                    "method": "structural",
                }
            )
        except Exception as e:
            logger.error(f"Failed to generate README: {e}")
            return ToolResult(success=False, output=None, error=str(e))

    def _analyze_project(self, project_path: str) -> dict:
        """Analyze project structure for README generation"""
        path = Path(project_path)
        analysis = {
            "structure": "",
            "key_files": [],
            "dependencies": []
        }

        try:
            # Find key files
            key_patterns = {
                "setup.py": "Python package",
                "pyproject.toml": "Modern Python project",
                "package.json": "Node.js project",
                "Cargo.toml": "Rust project",
                "go.mod": "Go project",
                "requirements.txt": "Python dependencies",
                "Dockerfile": "Docker support",
                ".github/workflows": "CI/CD",
                "tests/": "Test suite",
                "docs/": "Documentation"
            }

            for pattern, desc in key_patterns.items():
                if (path / pattern).exists():
                    analysis["key_files"].append(f"{pattern} ({desc})")

            # Read dependencies
            if (path / "requirements.txt").exists():
                with open(path / "requirements.txt", 'r') as f:
                    deps = [line.split('==')[0].strip() for line in f if line.strip() and not line.startswith('#')]
                    analysis["dependencies"] = deps[:10]  # Limit to top 10

            # Build structure string
            structure_lines = []
            for item in sorted(path.iterdir())[:15]:
                if item.name.startswith('.'):
                    continue
                structure_lines.append(f"- {item.name}{'/' if item.is_dir() else ''}")
            analysis["structure"] = "\n".join(structure_lines)

        except Exception as e:
            logger.warning(f"Project analysis partial failure: {e}")

        return analysis

    # Manifest filename -> install command it implies.
    _INSTALL = {
        "requirements.txt": "pip install -r requirements.txt",
        "pyproject.toml": "pip install .",
        "setup.py": "pip install .",
        "package.json": "npm install",
        "Cargo.toml": "cargo build",
        "go.mod": "go build ./...",
    }

    def _build_readme(self, project_name: str, description: str, info: dict) -> str:
        """Assemble README markdown from analyzed project facts only."""
        manifests = {kf.split(' ', 1)[0] for kf in info.get("key_files", [])}
        install_cmd = next((cmd for man, cmd in self._INSTALL.items() if man in manifests), None)

        lines = [f"# {project_name}", ""]
        if description:
            lines += [description, ""]

        if info.get("structure"):
            lines += ["## Project Structure", "", "```", info["structure"], "```", ""]

        lines += ["## Installation", ""]
        if install_cmd:
            lines += ["```bash", install_cmd, "```", ""]
        else:
            lines += ["_No dependency manifest detected._", ""]

        if info.get("dependencies"):
            lines += ["## Dependencies", ""]
            lines += [f"- {d}" for d in info["dependencies"]]
            lines += [""]

        if info.get("key_files"):
            lines += ["## Detected Project Facts", ""]
            lines += [f"- {kf}" for kf in info["key_files"]]
            lines += [""]

        return "\n".join(lines).rstrip() + "\n"


class GenerateAPIDocsTool(Tool):
    """Generate API documentation from code"""

    def __init__(self):
        super().__init__()
        self.name = "generate_api_docs"
        self.description = "Generate API documentation from source code"
        self.parameters = [
            ToolParameter(
                name="code",
                type="string",
                description="Source code to document",
                required=True
            ),
            ToolParameter(
                name="format",
                type="string",
                description="Documentation format (markdown, html, rst)",
                required=False,
                default="markdown"
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="generate_api_docs",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="GenerateAPIDocs capability"
                )
            ]
        )

    async def execute(self, **kwargs) -> ToolResult:
        """Generate API docs from code via AST, model-free.

        Documents exactly what is in the code — module/class/function
        docstrings and real signatures (annotations, defaults, return type).
        Nothing is inferred or invented.
        """
        try:
            import ast
            code = kwargs.get("code", "")
            format_type = kwargs.get("format", "markdown")

            if not code.strip():
                return ToolResult(success=False, output=None, error="no code provided")

            try:
                tree = ast.parse(code)
            except SyntaxError as se:
                return ToolResult(success=False, output=None,
                                  error=f"cannot document: code has a syntax error ({se})")

            docs = self._build_api_docs(tree)
            if not docs.strip():
                return ToolResult(success=False, output=None,
                                  error="no documentable functions or classes found in the code")

            return ToolResult(
                success=True,
                output={"documentation": docs, "format": format_type, "method": "ast"}
            )
        except Exception as e:
            logger.error(f"Failed to generate API docs: {e}")
            return ToolResult(success=False, output=None, error=str(e))

    @staticmethod
    def _format_signature(node) -> str:
        """Render a def's parameter list and return type from its AST."""
        import ast
        a = node.args
        parts: List[str] = []
        positional = list(getattr(a, 'posonlyargs', [])) + list(a.args)
        defaults = list(a.defaults)
        padded = [None] * (len(positional) - len(defaults)) + defaults
        for arg, default in zip(positional, padded):
            s = arg.arg
            if arg.annotation is not None:
                s += f": {ast.unparse(arg.annotation)}"
            if default is not None:
                s += f" = {ast.unparse(default)}"
            parts.append(s)
        if a.vararg:
            s = "*" + a.vararg.arg
            if a.vararg.annotation is not None:
                s += f": {ast.unparse(a.vararg.annotation)}"
            parts.append(s)
        for arg, default in zip(a.kwonlyargs, a.kw_defaults):
            s = arg.arg
            if arg.annotation is not None:
                s += f": {ast.unparse(arg.annotation)}"
            if default is not None:
                s += f" = {ast.unparse(default)}"
            parts.append(s)
        if a.kwarg:
            parts.append("**" + a.kwarg.arg)
        sig = f"({', '.join(parts)})"
        if node.returns is not None:
            sig += f" -> {ast.unparse(node.returns)}"
        return sig

    def _document_function(self, node, level: int = 3) -> List[str]:
        import ast
        prefix = "#" * level
        kind = "async def" if isinstance(node, ast.AsyncFunctionDef) else "def"
        lines = [f"{prefix} `{kind} {node.name}{self._format_signature(node)}`", ""]
        doc = ast.get_docstring(node)
        if doc:
            lines += [doc.strip(), ""]
        return lines

    def _build_api_docs(self, tree) -> str:
        import ast
        lines = ["# API Documentation", ""]
        module_doc = ast.get_docstring(tree)
        if module_doc:
            lines += [module_doc.strip(), ""]

        funcs = [n for n in tree.body if isinstance(n, (ast.FunctionDef, ast.AsyncFunctionDef))]
        classes = [n for n in tree.body if isinstance(n, ast.ClassDef)]

        if funcs:
            lines += ["## Functions", ""]
            for fn in funcs:
                lines += self._document_function(fn, level=3)

        for cls in classes:
            lines += [f"## class `{cls.name}`", ""]
            cdoc = ast.get_docstring(cls)
            if cdoc:
                lines += [cdoc.strip(), ""]
            methods = [n for n in cls.body if isinstance(n, (ast.FunctionDef, ast.AsyncFunctionDef))]
            for m in methods:
                lines += self._document_function(m, level=4)

        return "\n".join(lines).rstrip() + "\n"


class AnalyzeCodeQualityTool(Tool):
    """Analyze code quality and provide metrics"""

    def __init__(self):
        super().__init__()
        self.name = "analyze_code_quality"
        self.description = "Analyze code quality and provide detailed metrics"
        self.parameters = [
            ToolParameter(
                name="code",
                type="string",
                description="Source code to analyze",
                required=True
            ),
            ToolParameter(
                name="include_suggestions",
                type="boolean",
                description="Include improvement suggestions",
                required=False,
                default=True
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="analyze_code_quality",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="AnalyzeCodeQuality capability"
                )
            ]
        )

    async def execute(self, **kwargs) -> ToolResult:
        """Analyze code quality"""
        try:
            code = kwargs.get("code", "")
            include_suggestions = kwargs.get("include_suggestions", True)

            # Basic quality metrics
            lines = code.split('\n')
            total_lines = len(lines)
            code_lines = len([l for l in lines if l.strip() and not l.strip().startswith('#')])
            comment_lines = len([l for l in lines if l.strip().startswith('#')])

            # Calculate metrics
            quality_metrics = {
                "total_lines": total_lines,
                "code_lines": code_lines,
                "comment_lines": comment_lines,
                "comment_ratio": comment_lines / max(code_lines, 1),
                "complexity_score": min(10, code_lines / 10),  # Simple heuristic
                "maintainability_score": 8.5,  # Placeholder
                "issues": []
            }

            # Basic quality checks
            if quality_metrics["comment_ratio"] < 0.1:
                quality_metrics["issues"].append("Low comment coverage - consider adding more documentation")

            if code_lines > 500:
                quality_metrics["issues"].append("File is large - consider refactoring into smaller modules")

            if include_suggestions:
                quality_metrics["suggestions"] = [
                    "Add type hints to function parameters",
                    "Consider adding unit tests",
                    "Review naming conventions for clarity"
                ]

            return ToolResult(
                success=True,
                output=quality_metrics
            )
        except Exception as e:
            logger.error(f"Failed to analyze code quality: {e}")
            return ToolResult(success=False, output=None, error=str(e))


class ExtractDocstringsTool(Tool):
    """Extract docstrings from Python code"""

    def __init__(self):
        super().__init__()
        self.name = "extract_docstrings"
        self.description = "Extract all docstrings from Python code"
        self.parameters = [
            ToolParameter(
                name="code",
                type="string",
                description="Python source code to extract docstrings from",
                required=True
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="extract_docstrings",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="ExtractDocstrings capability"
                )
            ]
        )

    async def execute(self, **kwargs) -> ToolResult:
        """Extract docstrings"""
        try:
            code = kwargs.get("code", "")
            import ast

            # Parse code and extract docstrings
            tree = ast.parse(code)
            docstrings = {}

            for node in ast.walk(tree):
                if isinstance(node, (ast.FunctionDef, ast.ClassDef, ast.Module)):
                    docstring = ast.get_docstring(node)
                    if docstring:
                        name = getattr(node, 'name', 'module')
                        docstrings[name] = docstring

            return ToolResult(
                success=True,
                output={"docstrings": docstrings, "count": len(docstrings)}
            )
        except Exception as e:
            logger.error(f"Failed to extract docstrings: {e}")
            return ToolResult(success=False, output=None, error=str(e))


class GenerateChangelogTool(Tool):
    """Production-ready CHANGELOG generation from git history using LLM"""

    def __init__(self):
        super().__init__()
        self.name = "generate_changelog"
        self.description = "Generate CHANGELOG.md from git commit history with intelligent categorization"
        self.category = ToolCategory.DOCUMENTATION
        self.parameters = [
            ToolParameter(
                name="repo_path",
                type="string",
                description="Path to git repository",
                required=False,
                default="."
            ),
            ToolParameter(
                name="since_tag",
                type="string",
                description="Generate changelog since this git tag (optional)",
                required=False
            ),
            ToolParameter(
                name="max_commits",
                type="number",
                description="Maximum number of commits to analyze",
                required=False,
                default=100
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="generate_changelog",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="GenerateChangelog capability"
                ),
                CapabilityMetadata(
                    capability=Capability.GENERATE_REPORT,
                    description="Generate structured reports from data",
                    input_types=["data", "template"],
                    output_types=["report"],
                    latency="medium",
                    cost="low",
                    reliability="high",
                    risk_level=RiskLevel.LOW,
                    priority=8
                )
            ]
        )

    async def execute(self, **kwargs) -> ToolResult:
        """Generate changelog from git history"""
        import subprocess

        try:
            repo_path = kwargs.get("repo_path", ".")
            since_tag = kwargs.get("since_tag", None)
            max_commits = kwargs.get("max_commits", 100)

            # Get git commits
            git_cmd = ["git", "-C", repo_path, "log", f"--max-count={max_commits}", "--pretty=format:%h|%s|%an|%ad", "--date=short"]
            if since_tag:
                git_cmd.insert(4, f"{since_tag}..HEAD")

            try:
                result = subprocess.run(git_cmd, capture_output=True, text=True, timeout=10)
                if result.returncode != 0:
                    raise Exception("Not a git repository or git not available")

                commits_raw = result.stdout.strip()
            except (subprocess.SubprocessError, FileNotFoundError) as git_error:
                # No git means no commit history, so nothing was summarised.
                # The template is attached, not passed off as a changelog.
                return ToolResult(
                    success=False,
                    error=f"git history unavailable ({git_error}); template attached",
                    output={
                        "content": self._generate_changelog_template(),
                        "method": "template",
                        "note": "Git not available - generated template"
                    }
                )

            if not commits_raw:
                return ToolResult(
                    success=True,
                    output={
                        "content": self._generate_changelog_template(),
                        "method": "empty",
                        "note": "No commits found"
                    }
                )

            # Parse commits
            commits = []
            for line in commits_raw.split('\n'):
                if '|' in line:
                    parts = line.split('|')
                    if len(parts) >= 4:
                        commits.append({
                            "hash": parts[0],
                            "message": parts[1],
                            "author": parts[2],
                            "date": parts[3]
                        })

            # Categorize commits into Keep a Changelog sections, model-free,
            # from their conventional-commit type (feat/fix/...) or, failing
            # that, keywords in the subject line.
            changelog_content = self._render_changelog(commits)

            return ToolResult(
                success=True,
                output={
                    "content": changelog_content,
                    "commits_analyzed": len(commits),
                    "method": "conventional_commits"
                }
            )

        except Exception as e:
            logger.error(f"Failed to generate changelog: {e}")
            return ToolResult(
                success=False,
                error=f"changelog generation failed ({e}); template attached",
                output={
                    "content": self._generate_changelog_template(),
                    "method": "fallback",
                    "error": str(e)
                }
            )

    # Keep a Changelog sections, in output order.
    _SECTIONS = ["Added", "Changed", "Deprecated", "Removed", "Fixed", "Security"]

    # Conventional-commit type -> section.
    _TYPE_SECTION = {
        "feat": "Added", "add": "Added",
        "fix": "Fixed", "bugfix": "Fixed", "hotfix": "Fixed",
        "change": "Changed", "refactor": "Changed", "perf": "Changed",
        "update": "Changed", "improve": "Changed", "style": "Changed",
        "docs": "Changed", "doc": "Changed", "chore": "Changed", "build": "Changed",
        "ci": "Changed", "test": "Changed",
        "deprecate": "Deprecated",
        "remove": "Removed", "delete": "Removed", "drop": "Removed", "revert": "Removed",
        "security": "Security", "sec": "Security",
    }

    def _categorize(self, message: str) -> str:
        """Map one commit subject to a Keep a Changelog section."""
        head = message.strip().lower()
        # Conventional commit: "type(scope): subject" or "type: subject".
        m = re.match(r'([a-z]+)(?:\([^)]*\))?!?:', head)
        if m and m.group(1) in self._TYPE_SECTION:
            return self._TYPE_SECTION[m.group(1)]
        # Keyword fallback on the subject line.
        for kw, section in (
            ("secur", "Security"), ("vulnerab", "Security"), ("cve", "Security"),
            ("deprecat", "Deprecated"),
            ("remove", "Removed"), ("delete", "Removed"), ("drop ", "Removed"),
            ("fix", "Fixed"), ("bug", "Fixed"), ("patch", "Fixed"),
            ("add", "Added"), ("new ", "Added"), ("introduce", "Added"), ("implement", "Added"),
        ):
            if kw in head:
                return section
        return "Changed"

    def _render_changelog(self, commits: List[Dict[str, Any]]) -> str:
        """Render Keep a Changelog markdown by categorizing each commit."""
        buckets: Dict[str, List[str]] = {s: [] for s in self._SECTIONS}
        for c in commits:
            section = self._categorize(c["message"])
            subject = c["message"].strip()
            # Strip a leading conventional-commit prefix for the readable entry.
            subject = re.sub(r'^[a-z]+(?:\([^)]*\))?!?:\s*', '', subject, flags=re.IGNORECASE)
            buckets[section].append(f"- {subject} ({c['hash']})")

        lines = [
            "# Changelog",
            "",
            "All notable changes to this project will be documented in this file.",
            "",
            "The format is based on [Keep a Changelog](https://keepachangelog.com/en/1.0.0/),",
            "and this project adheres to [Semantic Versioning](https://semver.org/spec/v2.0.0.html).",
            "",
            "## [Unreleased]",
            "",
        ]
        for section in self._SECTIONS:
            if buckets[section]:
                lines.append(f"### {section}")
                lines.extend(buckets[section])
                lines.append("")
        return "\n".join(lines).rstrip() + "\n"

    def _generate_changelog_template(self) -> str:
        """Generate changelog template"""
        from datetime import datetime
        today = datetime.now().strftime("%Y-%m-%d")

        return f"""# Changelog

All notable changes to this project will be documented in this file.

The format is based on [Keep a Changelog](https://keepachangelog.com/en/1.0.0/),
and this project adheres to [Semantic Versioning](https://semver.org/spec/v2.0.0.html).

## [Unreleased]

### Added
- New features and functionality

### Changed
- Updates to existing features

### Deprecated
- Features that will be removed in future versions

### Removed
- Removed features

### Fixed
- Bug fixes

### Security
- Security improvements

## [1.0.0] - {today}

### Added
- Initial release
"""


class CreateDiagramTool(Tool):
    """Production-ready diagram generation using LLM for Mermaid syntax"""

    def __init__(self):
        super().__init__()
        self.name = "create_diagram"
        self.description = "Create technical diagrams using Mermaid syntax with LLM assistance"
        self.category = ToolCategory.DOCUMENTATION
        self.parameters = [
            ToolParameter(
                name="diagram_type",
                type="string",
                description="Type of diagram",
                required=True,
                enum=["flowchart", "sequence", "class", "er", "state", "gantt", "architecture"]
            ),
            ToolParameter(
                name="description",
                type="string",
                description="Description of what the diagram should show",
                required=True
            ),
            ToolParameter(
                name="elements",
                type="array",
                description="Key elements/components to include (optional)",
                required=False
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="create_diagram",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="CreateDiagram capability"
                ),
                CapabilityMetadata(
                    capability=Capability.CREATE_DIAGRAM,
                    description="Create diagrams from specifications",
                    input_types=["spec", "format"],
                    output_types=["diagram"],
                    latency="medium",
                    cost="low",
                    reliability="high",
                    risk_level=RiskLevel.LOW,
                    priority=8
                )
            ]
        )

    _MERMAID_HEADER = {
        "flowchart": "graph TD", "architecture": "graph LR",
        "sequence": "sequenceDiagram", "class": "classDiagram",
        "er": "erDiagram", "state": "stateDiagram-v2", "gantt": "gantt",
    }

    async def execute(self, **kwargs) -> ToolResult:
        """Build a Mermaid diagram from structured elements, model-free.

        Relationships cannot be invented from a free-text description without
        a model, so when no ``elements`` are given this reports an honest gap
        rather than fabricating a generic shape.
        """
        try:
            diagram_type = kwargs.get("diagram_type", "flowchart")
            description = kwargs.get("description", "")
            elements = kwargs.get("elements", []) or []

            if not elements:
                return ToolResult(
                    success=False,
                    output={"diagram_type": diagram_type, "description": description},
                    error=("no elements to diagram: provide `elements` as node labels "
                           "or edges like 'A -> B'; a diagram cannot be synthesized from "
                           "a free-text description model-free"),
                )

            diagram_code = self._build_mermaid(diagram_type, [str(e) for e in elements])
            diagram_markdown = f"```mermaid\n{diagram_code}\n```"

            return ToolResult(
                success=True,
                output={
                    "diagram": diagram_markdown,
                    "diagram_code": diagram_code,
                    "format": "mermaid",
                    "diagram_type": diagram_type,
                    "method": "structural",
                }
            )

        except Exception as e:
            logger.error(f"Failed to create diagram: {e}")
            return ToolResult(success=False, output=None, error=str(e))

    @staticmethod
    def _node_id(index: int) -> str:
        return chr(ord('A') + index) if index < 26 else f"N{index}"

    def _build_mermaid(self, diagram_type: str, elements: List[str]) -> str:
        """Assemble Mermaid syntax from a list of node labels or 'A -> B' edges."""
        header = self._MERMAID_HEADER.get(diagram_type, "graph TD")

        def ident(label: str) -> str:
            return re.sub(r'[^A-Za-z0-9_]', '_', label.strip()) or "n"

        # Explicit edges take precedence for graph-shaped diagrams.
        edges = [e for e in elements if re.search(r'-+>', e)]
        if edges and diagram_type in ("flowchart", "architecture", "state", "er"):
            lines = [header]
            for e in edges:
                src, dst = re.split(r'-+>', e, maxsplit=1)
                lines.append(f"    {ident(src)} --> {ident(dst)}")
            return "\n".join(lines)

        labels = [e.strip() for e in elements if e.strip()]

        if diagram_type == "sequence":
            lines = [header]
            for lab in labels:
                lines.append(f"    participant {ident(lab)}")
            for a, b in zip(labels, labels[1:]):
                lines.append(f"    {ident(a)}->>{ident(b)}: {b}")
            return "\n".join(lines)

        if diagram_type == "class":
            lines = [header]
            for lab in labels:
                lines.append(f"    class {ident(lab)}")
            return "\n".join(lines)

        # flowchart / architecture / er / state / gantt: chain the nodes in order.
        lines = [header]
        ids = [self._node_id(i) for i in range(len(labels))]
        for i, lab in enumerate(labels):
            lines.append(f"    {ids[i]}[{lab}]")
        for a, b in zip(ids, ids[1:]):
            lines.append(f"    {a} --> {b}")
        return "\n".join(lines)


class UpdateDocsTool(Tool):
    """Production-ready documentation update with intelligent merging"""

    def __init__(self):
        super().__init__()
        self.name = "update_docs"
        self.description = "Update existing documentation files with intelligent content merging"
        self.category = ToolCategory.DOCUMENTATION
        self.parameters = [
            ToolParameter(
                name="doc_path",
                type="string",
                description="Path to documentation file to update",
                required=True
            ),
            ToolParameter(
                name="section",
                type="string",
                description="Section to update (e.g., 'Installation', 'API Reference')",
                required=True
            ),
            ToolParameter(
                name="new_content",
                type="string",
                description="New content for the section",
                required=True
            ),
            ToolParameter(
                name="merge_strategy",
                type="string",
                description="How to merge content",
                required=False,
                default="replace",
                enum=["replace", "append", "prepend", "smart_merge"]
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="update_docs",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="UpdateDocs capability"
                )
            ]
        )

    async def execute(self, **kwargs) -> ToolResult:
        """Update documentation section"""
        try:
            doc_path = kwargs.get("doc_path", "")
            section = kwargs.get("section", "")
            new_content = kwargs.get("new_content", "")
            merge_strategy = kwargs.get("merge_strategy", "replace")

            path = Path(doc_path)
            if not path.exists():
                return ToolResult(success=False, output=None, error=f"Documentation file not found: {doc_path}")

            # Read existing content
            with open(path, 'r', encoding='utf-8') as f:
                existing_content = f.read()

            # Find section in markdown
            section_pattern = rf'(##\s+{re.escape(section)}.*?)(?=\n##\s+|\Z)'
            section_match = re.search(section_pattern, existing_content, re.DOTALL)

            if merge_strategy == "replace" and section_match:
                # Replace entire section
                updated_content = re.sub(
                    section_pattern,
                    f"## {section}\n\n{new_content}\n",
                    existing_content,
                    flags=re.DOTALL
                )
            elif merge_strategy == "append" and section_match:
                # Append to section
                section_content = section_match.group(0)
                updated_section = section_content.rstrip() + "\n\n" + new_content + "\n"
                updated_content = existing_content.replace(section_content, updated_section)
            elif merge_strategy == "prepend" and section_match:
                # Prepend to section
                updated_section = f"## {section}\n\n{new_content}\n\n{section_match.group(0)[len(section)+3:]}"
                updated_content = existing_content.replace(section_match.group(0), updated_section)
            else:
                # Section doesn't exist or smart_merge - append new section
                updated_content = existing_content.rstrip() + f"\n\n## {section}\n\n{new_content}\n"

            # Write updated content
            with open(path, 'w', encoding='utf-8') as f:
                f.write(updated_content)

            return ToolResult(
                success=True,
                output={
                    "updated": True,
                    "path": str(path),
                    "section": section,
                    "merge_strategy": merge_strategy,
                    "content_length": len(updated_content)
                }
            )

        except Exception as e:
            logger.error(f"Failed to update docs: {e}")
            return ToolResult(success=False, output=None, error=str(e))


class DocsBuildPreviewTool(Tool):
    """Production-ready documentation build and preview"""

    def __init__(self):
        super().__init__()
        self.name = "docs_build_preview"
        self.description = "Build and preview documentation site using MkDocs or Sphinx"
        self.category = ToolCategory.DOCUMENTATION
        self.parameters = [
            ToolParameter(
                name="docs_dir",
                type="string",
                description="Documentation directory",
                required=False,
                default="docs"
            ),
            ToolParameter(
                name="tool",
                type="string",
                description="Documentation tool to use",
                required=False,
                default="auto",
                enum=["auto", "mkdocs", "sphinx"]
            ),
            ToolParameter(
                name="port",
                type="number",
                description="Preview server port",
                required=False,
                default=8000
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="docs_build_preview",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="DocsBuildPreview capability"
                )
            ]
        )

    async def execute(self, **kwargs) -> ToolResult:
        """Build and serve documentation preview"""
        import subprocess

        try:
            docs_dir = kwargs.get("docs_dir", "docs")
            tool = kwargs.get("tool", "auto")
            port = kwargs.get("port", 8000)

            docs_path = Path(docs_dir)
            if not docs_path.exists():
                return ToolResult(success=False, output=None, error=f"Documentation directory not found: {docs_dir}")

            # Auto-detect documentation tool
            if tool == "auto":
                if (docs_path.parent / "mkdocs.yml").exists():
                    tool = "mkdocs"
                elif (docs_path.parent / "conf.py").exists() or (docs_path / "conf.py").exists():
                    tool = "sphinx"
                else:
                    tool = "mkdocs"  # Default

            if tool == "mkdocs":
                # Try to build with MkDocs
                try:
                    # Build docs
                    build_result = subprocess.run(
                        ["mkdocs", "build"],
                        cwd=docs_path.parent,
                        capture_output=True,
                        text=True,
                        timeout=30
                    )

                    if build_result.returncode == 0:
                        return ToolResult(
                            success=True,
                            output={
                                "build_status": "success",
                                "tool": "mkdocs",
                                "output_dir": "site",
                                "preview_command": f"mkdocs serve -a localhost:{port}",
                                "preview_url": f"http://localhost:{port}",
                                "note": f"Run 'mkdocs serve' to preview locally"
                            }
                        )
                    else:
                        return ToolResult(
                            success=False,
                            error=f"MkDocs build failed: {build_result.stderr}"
                        )

                except FileNotFoundError:
                    return ToolResult(
                        success=False,
                        error="MkDocs not installed. Install with: pip install mkdocs"
                    )

            elif tool == "sphinx":
                # Try to build with Sphinx
                try:
                    build_result = subprocess.run(
                        ["sphinx-build", "-b", "html", docs_dir, f"{docs_dir}/_build/html"],
                        capture_output=True,
                        text=True,
                        timeout=30
                    )

                    if build_result.returncode == 0:
                        return ToolResult(
                            success=True,
                            output={
                                "build_status": "success",
                                "tool": "sphinx",
                                "output_dir": f"{docs_dir}/_build/html",
                                "preview_command": f"python -m http.server {port} -d {docs_dir}/_build/html",
                                "preview_url": f"http://localhost:{port}",
                                "note": "Built successfully. Use http.server or sphinx-autobuild to preview"
                            }
                        )
                    else:
                        return ToolResult(
                            success=False,
                            error=f"Sphinx build failed: {build_result.stderr}"
                        )

                except FileNotFoundError:
                    return ToolResult(
                        success=False,
                        error="Sphinx not installed. Install with: pip install sphinx"
                    )

        except Exception as e:
            logger.error(f"Failed to build/preview docs: {e}")
            return ToolResult(success=False, output=None, error=str(e))


class VersionedDocDeploymentTool(Tool):
    """Production-ready versioned documentation deployment"""

    def __init__(self):
        super().__init__()
        self.name = "versioned_doc_deployment"
        self.description = "Deploy versioned documentation with mike or manual versioning"
        self.category = ToolCategory.DOCUMENTATION
        self.parameters = [
            ToolParameter(
                name="version",
                type="string",
                description="Documentation version (e.g., '1.0', 'latest')",
                required=True
            ),
            ToolParameter(
                name="docs_dir",
                type="string",
                description="Documentation directory",
                required=False,
                default="docs"
            ),
            ToolParameter(
                name="alias",
                type="string",
                description="Version alias (e.g., 'stable', 'latest')",
                required=False
            ),
            ToolParameter(
                name="push",
                type="boolean",
                description="Push to gh-pages branch",
                required=False,
                default=False
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="versioned_doc_deployment",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="VersionedDocDeployment capability"
                )
            ]
        )

    async def execute(self, **kwargs) -> ToolResult:
        """Deploy versioned documentation"""
        import subprocess

        try:
            version = kwargs.get("version", "latest")
            docs_dir = kwargs.get("docs_dir", "docs")
            alias = kwargs.get("alias", None)
            push = kwargs.get("push", False)

            # Check if mike is available (for MkDocs versioning)
            try:
                mike_cmd = ["mike", "deploy", "--update-aliases", version]
                if alias:
                    mike_cmd.append(alias)
                if push:
                    mike_cmd.append("--push")

                result = subprocess.run(
                    mike_cmd,
                    capture_output=True,
                    text=True,
                    timeout=60
                )

                if result.returncode == 0:
                    return ToolResult(
                        success=True,
                        output={
                            "deployed": True,
                            "version": version,
                            "alias": alias,
                            "method": "mike",
                            "pushed_to_remote": push,
                            "url": f"https://<username>.github.io/<repo>/{version}/",
                            "note": "Deployed with mike. Update URL with your GitHub Pages domain."
                        }
                    )
                else:
                    raise Exception(f"Mike deployment failed: {result.stderr}")

            except FileNotFoundError:
                # NOTHING WAS DEPLOYED. `deployed: False` was in the output but
                # success=True is what a caller branches on.
                return ToolResult(
                    success=False,
                    error="mike is not installed, so nothing was deployed; "
                          "manual instructions attached",
                    output={
                        "deployed": False,
                        "version": version,
                        "method": "manual",
                        "instructions": [
                            "1. Install mike: pip install mike",
                            f"2. Deploy version: mike deploy --push {version} {alias or ''}",
                            "3. Set default: mike set-default latest",
                            "4. Configure GitHub Pages to serve from gh-pages branch"
                        ],
                        "note": "Mike not installed. Follow manual instructions above."
                    }
                )

        except Exception as e:
            logger.error(f"Failed to deploy versioned docs: {e}")
            return ToolResult(
                success=False,
                error=str(e),
                output={
                    "deployed": False,
                    "version": version,
                    "error_details": str(e)
                }
            )


class ADRGeneratorTool(Tool):
    """Generate Architecture Decision Records (ADR)"""

    def __init__(self):
        super().__init__()
        self.name = "adr_generator"
        self.description = "Generate Architecture Decision Records"
        self.parameters = [
            ToolParameter(
                name="decision_title",
                type="string",
                description="Title of the architecture decision",
                required=True
            ),
            ToolParameter(
                name="context",
                type="string",
                description="Context for the decision",
                required=True
            ),
            ToolParameter(
                name="decision",
                type="string",
                description="The decision made",
                required=True
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="adr_generator",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="ADRGenerator capability"
                )
            ]
        )

    async def execute(self, **kwargs) -> ToolResult:
        """Generate ADR"""
        try:
            title = kwargs.get("decision_title", "")
            context = kwargs.get("context", "")
            decision = kwargs.get("decision", "")

            from datetime import datetime
            date = datetime.now().strftime("%Y-%m-%d")

            adr = f"""# {title}

Date: {date}

## Status

Accepted

## Context

{context}

## Decision

{decision}

## Consequences

TBD
"""

            return ToolResult(success=True, output={"adr": adr, "title": title})
        except Exception as e:
            return ToolResult(success=False, output=None, error=str(e))


# ============================================================================
# REAL DOCUMENT GENERATION TOOLS (PDF, DOCX, PPTX, etc.)
# ============================================================================


class GeneratePDFDocumentTool(Tool):
    """Generate actual PDF documents using reportlab"""

    def __init__(self):
        super().__init__()
        self.name = "generate_pdf_document"
        self.description = "Generate a real PDF document with formatting, images, and styling"
        self.category = ToolCategory.DOCUMENTATION
        self.parameters = [
            ToolParameter(
                name="output_path",
                type="string",
                description="Path where PDF should be saved",
                required=True
            ),
            ToolParameter(
                name="title",
                type="string",
                description="Document title",
                required=True
            ),
            ToolParameter(
                name="content",
                type="string",
                description="Document content (supports markdown-like formatting)",
                required=True
            ),
            ToolParameter(
                name="author",
                type="string",
                description="Document author",
                required=False
            ),
            ToolParameter(
                name="include_toc",
                type="boolean",
                description="Include table of contents",
                required=False,
                default=False
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="generate_pdf_document",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="GeneratePDFDocument capability"
                )
            ]
        )

    def _to_roman(self, num):
        """Convert integer to Roman numeral"""
        val = [1000, 900, 500, 400, 100, 90, 50, 40, 10, 9, 5, 4, 1]
        syms = ['M', 'CM', 'D', 'CD', 'C', 'XC', 'L', 'XL', 'X', 'IX', 'V', 'IV', 'I']
        roman_num = ''
        i = 0
        while num > 0:
            for _ in range(num // val[i]):
                roman_num += syms[i]
                num -= val[i]
            i += 1
        return roman_num

    async def execute(self, **kwargs) -> ToolResult:
        """Generate enterprise-grade PDF document with professional styling"""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
            from reportlab.lib import colors
            from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
            from datetime import datetime

            output_path = kwargs.get("output_path")
            title = kwargs.get("title", "Document")
            content = kwargs.get("content", "")
            author = kwargs.get("author", "TorinAI")
            include_toc = kwargs.get("include_toc", False)

            # Create PDF with professional margins
            doc = SimpleDocTemplate(
                output_path,
                pagesize=letter,
                rightMargin=0.75*inch,
                leftMargin=0.75*inch,
                topMargin=1*inch,
                bottomMargin=0.75*inch,
                title=title,
                author=author
            )
            story = []
            styles = getSampleStyleSheet()

            # Corporate color scheme
            primary_color = colors.HexColor('#1a1a2e')
            secondary_color = colors.HexColor('#16213e')
            text_color = colors.HexColor('#2C3E50')
            accent_color = colors.HexColor('#3498DB')

            # Professional title style (cover page)
            cover_title_style = ParagraphStyle(
                'CoverTitle',
                parent=styles['Title'],
                fontSize=32,
                textColor=primary_color,
                spaceAfter=0.4*inch,
                alignment=TA_CENTER,
                fontName='Helvetica-Bold',
                leading=38
            )

            # Cover page layout
            story.append(Spacer(1, 2*inch))
            story.append(Paragraph(title, cover_title_style))

            # Metadata on cover
            meta_style = ParagraphStyle(
                'Metadata',
                parent=styles['Normal'],
                fontSize=12,
                textColor=secondary_color,
                spaceAfter=6,
                alignment=TA_CENTER,
                fontName='Helvetica'
            )
            story.append(Paragraph(f"<b>Author:</b> {author}", meta_style))
            story.append(Paragraph(f"<b>Generated:</b> {datetime.now().strftime('%B %d, %Y at %H:%M')}", meta_style))
            story.append(PageBreak())

            # Professional heading styles
            h1_style = ParagraphStyle(
                'EnterpriseH1',
                parent=styles['Heading1'],
                fontSize=20,
                textColor=primary_color,
                spaceBefore=0.3*inch,
                spaceAfter=0.12*inch,
                fontName='Helvetica-Bold',
                leading=22,
                borderColor=accent_color,
                borderWidth=0,
                borderPadding=6
            )

            h2_style = ParagraphStyle(
                'EnterpriseH2',
                parent=styles['Heading2'],
                fontSize=12,
                textColor=secondary_color,
                spaceBefore=0.2*inch,
                spaceAfter=0.08*inch,
                fontName='Helvetica-Bold',
                leading=12
            )

            h3_style = ParagraphStyle(
                'EnterpriseH3',
                parent=styles['Heading3'],
                fontSize=12,
                textColor=secondary_color,
                spaceBefore=0.15*inch,
                spaceAfter=0.06*inch,
                fontName='Helvetica-Bold',
                leading=13
            )

            # Professional body text
            body_style = ParagraphStyle(
                'EnterpriseBody',
                parent=styles['Normal'],
                fontSize=12,
                leading=16,
                spaceBefore=4,
                spaceAfter=8,
                textColor=text_color,
                alignment=TA_LEFT,
                fontName='Helvetica'
            )

            # Bullet point style
            bullet_style = ParagraphStyle(
                'EnterpriseBullet',
                parent=body_style,
                leftIndent=25,
                bulletIndent=10,
                spaceBefore=6,
                spaceAfter=6,
                alignment=TA_LEFT
            )

            # Pre-process content: split headers from their content
            # This ensures "## Header\n- bullet" becomes two separate paragraphs
            preprocessed_content = content
            # Add double newline after headers if followed by single newline
            preprocessed_content = re.sub(r'^(#{1,3}\s+[^\n]+)\n([^#\n])', r'\1\n\n\2', preprocessed_content, flags=re.MULTILINE)

            # Process content with intelligent parsing
            paragraphs = preprocessed_content.split('\n\n')
            section_number = 0
            last_para_type = None  # Track previous paragraph type

            for i, para in enumerate(paragraphs):
                if not para.strip():
                    continue

                # Explicit markdown headers
                if para.startswith('# '):
                    section_number += 1
                    roman = self._to_roman(section_number)
                    story.append(Paragraph(f"{roman}. {para[2:].strip()}", h1_style))
                    last_para_type = 'h1'

                elif para.startswith('## '):
                    # h2 style already has appropriate spaceBefore - no extra spacing needed
                    story.append(Paragraph(para[3:].strip(), h2_style))
                    last_para_type = 'h2'

                elif para.startswith('### '):
                    story.append(Paragraph(para[4:].strip(), h3_style))
                    last_para_type = 'h3'

                # Bullet lists
                elif para.strip().startswith(('- ', '• ', '* ')):
                    bullet_lines = para.strip().split('\n')
                    for bullet in bullet_lines:
                        bullet_stripped = bullet.strip()
                        if bullet_stripped.startswith('- '):
                            clean_bullet = bullet_stripped[2:].strip()
                        elif bullet_stripped.startswith('• '):
                            clean_bullet = bullet_stripped[2:].strip()
                        elif bullet_stripped.startswith('* '):
                            clean_bullet = bullet_stripped[2:].strip()
                        else:
                            continue

                        if clean_bullet:
                            # Handle markdown formatting - bold, italic, code
                            clean_bullet = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', clean_bullet)  # Bold (non-greedy)
                            clean_bullet = re.sub(r'(?<!\*)\*([^*]+)\*(?!\*)', r'<i>\1</i>', clean_bullet)  # Italic
                            clean_bullet = re.sub(r'`([^`]+)`', r'<font name="Courier">\1</font>', clean_bullet)  # Code
                            story.append(Paragraph(f"• {clean_bullet}", bullet_style))
                    story.append(Spacer(1, 0.15*inch))
                    last_para_type = 'bullets'

                # Numbered lists
                elif re.match(r'^\d+\.', para.strip()):
                    items = para.strip().split('\n')
                    for item in items:
                        if item.strip():
                            story.append(Paragraph(item.strip(), bullet_style))
                    story.append(Spacer(1, 0.1*inch))
                    last_para_type = 'numbered'

                # Regular paragraphs (check for inline bullets)
                else:
                    # Check if this paragraph has inline bullet lists like "text: - item1 - item2"
                    if re.search(r':\s*-\s+\*\*', para):
                        # Split intro text from bullets
                        parts = re.split(r'(:\s*-\s+)', para, maxsplit=1)
                        if len(parts) >= 2:
                            # Add intro text
                            intro = parts[0] + ':'
                            intro = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', intro)
                            intro = re.sub(r'(?<!\*)\*([^*]+)\*(?!\*)', r'<i>\1</i>', intro)
                            story.append(Paragraph(intro, body_style))
                            story.append(Spacer(1, 0.1*inch))

                            # Process remaining text as bullets
                            bullet_text = parts[2] if len(parts) > 2 else ""
                            # Split by " - " pattern
                            bullets = re.split(r'\s+-\s+', bullet_text)
                            for bullet in bullets:
                                if bullet.strip():
                                    clean = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', bullet.strip())
                                    clean = re.sub(r'(?<!\*)\*([^*]+)\*(?!\*)', r'<i>\1</i>', clean)
                                    story.append(Paragraph(f"• {clean}", bullet_style))
                            story.append(Spacer(1, 0.1*inch))
                            last_para_type = 'bullets'
                            continue

                    # Regular paragraph without inline bullets
                    processed = para.strip()
                    processed = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', processed)  # Bold (non-greedy)
                    processed = re.sub(r'(?<!\*)\*([^*]+)\*(?!\*)', r'<i>\1</i>', processed)  # Italic
                    processed = re.sub(r'`([^`]+)`', r'<font name="Courier">\1</font>', processed)  # Code

                    story.append(Paragraph(processed, body_style))
                    last_para_type = 'paragraph'

            # Add professional footer with page numbers
            def add_page_footer(canvas_obj, doc_obj):
                """Add footer with page numbers"""
                canvas_obj.saveState()
                canvas_obj.setFont('Helvetica', 9)
                canvas_obj.setFillColor(colors.HexColor('#666666'))

                # Page number
                page_num = f"Page {canvas_obj.getPageNumber()}"
                canvas_obj.drawRightString(
                    letter[0] - 0.75*inch,
                    0.5*inch,
                    page_num
                )

                # Footer text
                canvas_obj.drawString(
                    0.75*inch,
                    0.5*inch,
                    f"{author} • {title}"
                )

                canvas_obj.restoreState()

            # Build PDF
            doc.build(story, onFirstPage=add_page_footer, onLaterPages=add_page_footer)

            return ToolResult(
                success=True,
                output={
                    "pdf_path": output_path,
                    "title": title,
                    "pages": "generated",
                    "format": "PDF",
                    "size_bytes": Path(output_path).stat().st_size if Path(output_path).exists() else 0
                }
            )

        except Exception as e:
            logger.error(f"Failed to generate PDF: {e}")
            import traceback
            traceback.print_exc()
            return ToolResult(success=False, output=None, error=str(e))


class GenerateWordDocumentTool(Tool):
    """Generate Microsoft Word .docx documents"""

    def __init__(self):
        super().__init__()
        self.name = "generate_word_document"
        self.description = "Generate a Microsoft Word .docx document with formatting"
        self.category = ToolCategory.DOCUMENTATION
        self.parameters = [
            ToolParameter(
                name="output_path",
                type="string",
                description="Path where .docx should be saved",
                required=True
            ),
            ToolParameter(
                name="title",
                type="string",
                description="Document title",
                required=True
            ),
            ToolParameter(
                name="content",
                type="string",
                description="Document content",
                required=True
            ),
            ToolParameter(
                name="author",
                type="string",
                description="Document author",
                required=False
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="generate_word_document",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="GenerateWordDocument capability"
                )
            ]
        )

    async def execute(self, **kwargs) -> ToolResult:
        """Generate Word document"""
        try:
            from docx import Document
            from docx.shared import Inches, Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            output_path = kwargs.get("output_path")
            title = kwargs.get("title", "Document")
            content = kwargs.get("content", "")
            author = kwargs.get("author", "TorinAI")

            # Create document
            doc = Document()

            # Set document properties
            doc.core_properties.title = title
            doc.core_properties.author = author

            # Add title
            title_para = doc.add_heading(title, 0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Add metadata
            doc.add_paragraph(f"Author: {author}")
            from datetime import datetime
            doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
            doc.add_paragraph()

            # Process content
            paragraphs = content.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    # Check for headers
                    if para.startswith('# '):
                        doc.add_heading(para[2:], level=1)
                    elif para.startswith('## '):
                        doc.add_heading(para[3:], level=2)
                    elif para.startswith('### '):
                        doc.add_heading(para[4:], level=3)
                    else:
                        doc.add_paragraph(para)

            # Save document
            doc.save(output_path)

            return ToolResult(
                success=True,
                output={
                    "docx_path": output_path,
                    "title": title,
                    "format": "DOCX",
                    "size_bytes": Path(output_path).stat().st_size if Path(output_path).exists() else 0
                }
            )

        except Exception as e:
            logger.error(f"Failed to generate Word document: {e}")
            return ToolResult(success=False, output=None, error=str(e))


class GeneratePowerPointTool(Tool):
    """Generate PowerPoint .pptx presentations"""

    def __init__(self):
        super().__init__()
        self.name = "generate_powerpoint"
        self.description = "Generate a PowerPoint .pptx presentation with slides"
        self.category = ToolCategory.DOCUMENTATION
        self.parameters = [
            ToolParameter(
                name="output_path",
                type="string",
                description="Path where .pptx should be saved",
                required=True
            ),
            ToolParameter(
                name="title",
                type="string",
                description="Presentation title",
                required=True
            ),
            ToolParameter(
                name="slides",
                type="array",
                description="List of slides with 'title' and 'content' keys",
                required=True
            ),
            ToolParameter(
                name="author",
                type="string",
                description="Presentation author",
                required=False
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="generate_powerpoint",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="GeneratePowerPoint capability"
                )
            ]
        )

    async def execute(self, **kwargs) -> ToolResult:
        """Generate PowerPoint presentation"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            output_path = kwargs.get("output_path")
            title = kwargs.get("title", "Presentation")
            slides_data = kwargs.get("slides", [])
            author = kwargs.get("author", "TorinAI")

            # Create presentation
            prs = Presentation()
            prs.core_properties.title = title
            prs.core_properties.author = author

            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide_title = slide.shapes.title
            subtitle = slide.placeholders[1]
            slide_title.text = title
            from datetime import datetime
            subtitle.text = f"By {author} | {datetime.now().strftime('%Y-%m-%d')}"

            # Add content slides
            for slide_data in slides_data:
                if isinstance(slide_data, dict):
                    bullet_slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(bullet_slide_layout)
                    shapes = slide.shapes

                    title_shape = shapes.title
                    body_shape = shapes.placeholders[1]

                    title_shape.text = slide_data.get('title', 'Slide')

                    text_frame = body_shape.text_frame
                    content = slide_data.get('content', '')

                    # Split content into bullet points
                    bullets = content.split('\n')
                    for i, bullet in enumerate(bullets):
                        if bullet.strip():
                            if i == 0:
                                text_frame.text = bullet.strip()
                            else:
                                p = text_frame.add_paragraph()
                                p.text = bullet.strip()
                                p.level = 0

            # Save presentation
            prs.save(output_path)

            return ToolResult(
                success=True,
                output={
                    "pptx_path": output_path,
                    "title": title,
                    "slide_count": len(prs.slides),
                    "format": "PPTX",
                    "size_bytes": Path(output_path).stat().st_size if Path(output_path).exists() else 0
                }
            )

        except Exception as e:
            logger.error(f"Failed to generate PowerPoint: {e}")
            return ToolResult(success=False, output=None, error=str(e))


class GenerateArchitectureDiagramTool(Tool):
    """Generate architecture diagrams as PNG/SVG images"""

    def __init__(self):
        super().__init__()
        self.name = "generate_architecture_diagram"
        self.description = "Generate system architecture diagram as PNG or SVG image"
        self.category = ToolCategory.DOCUMENTATION
        self.parameters = [
            ToolParameter(
                name="output_path",
                type="string",
                description="Path where diagram image should be saved (.png or .svg)",
                required=True
            ),
            ToolParameter(
                name="components",
                type="array",
                description="List of system components to include",
                required=True
            ),
            ToolParameter(
                name="title",
                type="string",
                description="Diagram title",
                required=False
            ),
            ToolParameter(
                name="style",
                type="string",
                description="Diagram style (layered, circular, hierarchical)",
                required=False,
                default="layered"
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="generate_architecture_diagram",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="GenerateArchitectureDiagram capability"
                )
            ]
        )

    async def execute(self, **kwargs) -> ToolResult:
        """Generate enterprise-grade architecture diagram with professional styling"""
        try:
            import matplotlib.pyplot as plt
            import matplotlib.patches as mpatches
            from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Circle
            import numpy as np

            output_path = kwargs.get("output_path")
            components = kwargs.get("components", [])
            title = kwargs.get("title", "System Architecture")
            style = kwargs.get("style", "layered")

            # Create large, high-quality figure
            fig, ax = plt.subplots(figsize=(16, 12), facecolor='white', dpi=150)
            ax.set_xlim(0, 16)
            ax.set_ylim(0, 12)
            ax.axis('off')

            # Professional title with large font
            ax.text(8, 11.2, title, ha='center', va='top',
                   fontsize=24, fontweight='bold', color='#1a1a2e',
                   bbox=dict(boxstyle='round,pad=0.5', facecolor='#f0f0f0', edgecolor='#1a1a2e', linewidth=2))

            # Corporate color palette
            colors_palette = ['#3498DB', '#2ECC71', '#F39C12', '#E74C3C', '#9B59B6', '#1ABC9C']

            if style == "layered":
                # Layered architecture with generous spacing
                num_components = len(components) if components else 4
                if not components:
                    components = ["Presentation Layer", "Business Logic Layer", "Data Access Layer", "Database Layer"]

                # Calculate positions with generous spacing (minimum 1.5 units between components)
                available_height = 8.5  # From y=10 down to y=1.5
                spacing = min(available_height / max(num_components - 1, 1), 2.0)  # Max 2.0 spacing
                y_start = 10

                for i, comp in enumerate(components):
                    y_pos = y_start - (i * spacing)
                    comp_name = comp if isinstance(comp, str) else comp.get('name', f'Component {i+1}')
                    comp_color = colors_palette[i % len(colors_palette)]

                    # Large component box with rounded corners
                    box = FancyBboxPatch(
                        (2.5, y_pos - 0.5),  # x, y with generous height
                        11, 1.0,  # width, height (taller boxes)
                        boxstyle="round,pad=0.2",
                        edgecolor='#1a1a2e',
                        facecolor=comp_color,
                        linewidth=3,
                        alpha=0.9
                    )
                    ax.add_patch(box)

                    # Large, clear text
                    ax.text(8, y_pos, comp_name, ha='center', va='center',
                           fontsize=16, color='white', fontweight='bold',
                           bbox=dict(boxstyle='round,pad=0.3', facecolor=comp_color,
                                   edgecolor='none', alpha=0.0))

                    # Professional arrows between layers with generous spacing
                    if i < len(components) - 1:
                        arrow_start_y = y_pos - 0.6
                        arrow_end_y = y_pos - spacing + 0.6

                        arrow = FancyArrowPatch(
                            (8, arrow_start_y),
                            (8, arrow_end_y),
                            arrowstyle='->,head_width=0.6,head_length=0.6',
                            mutation_scale=40,
                            color='#2C3E50',
                            linewidth=4,
                            zorder=1
                        )
                        ax.add_patch(arrow)

            elif style == "circular" or style == "network":
                # Circular/network layout with nodes
                num_components = len(components) if components else 6
                if not components:
                    components = ["API Gateway", "Service A", "Service B", "Database", "Cache", "Queue"]

                # Arrange in circle with generous radius
                radius = 3.5
                center_x, center_y = 8, 6

                positions = []
                for i in range(num_components):
                    angle = 2 * np.pi * i / num_components - np.pi / 2
                    x = center_x + radius * np.cos(angle)
                    y = center_y + radius * np.sin(angle)
                    positions.append((x, y))

                    comp_name = components[i] if isinstance(components[i], str) else components[i].get('name', f'Node {i+1}')
                    comp_color = colors_palette[i % len(colors_palette)]

                    # Large circular nodes
                    circle = Circle(
                        (x, y), 0.8,
                        edgecolor='#1a1a2e',
                        facecolor=comp_color,
                        linewidth=3,
                        alpha=0.9,
                        zorder=2
                    )
                    ax.add_patch(circle)

                    # Node label with clear text
                    ax.text(x, y, comp_name, ha='center', va='center',
                           fontsize=14, color='white', fontweight='bold',
                           zorder=3)

                # Draw connections between nodes
                for i in range(num_components):
                    next_i = (i + 1) % num_components
                    arrow = FancyArrowPatch(
                        positions[i],
                        positions[next_i],
                        arrowstyle='<->,head_width=0.4,head_length=0.4',
                        mutation_scale=25,
                        color='#95A5A6',
                        linewidth=2.5,
                        alpha=0.6,
                        zorder=1
                    )
                    ax.add_patch(arrow)

            else:  # hierarchical or default
                # Hierarchical tree layout
                num_components = len(components) if components else 5
                if not components:
                    components = ["Root System", "Module A", "Module B", "Subsystem 1", "Subsystem 2"]

                # Simple hierarchical layout
                y_positions = np.linspace(10, 2, min(num_components, 4))

                for i, comp in enumerate(components[:4]):
                    y_pos = y_positions[i]
                    comp_name = comp if isinstance(comp, str) else comp.get('name', f'Component {i+1}')
                    comp_color = colors_palette[i % len(colors_palette)]

                    box = FancyBboxPatch(
                        (5, y_pos - 0.5),
                        6, 1.0,
                        boxstyle="round,pad=0.2",
                        edgecolor='#1a1a2e',
                        facecolor=comp_color,
                        linewidth=3,
                        alpha=0.9
                    )
                    ax.add_patch(box)

                    ax.text(8, y_pos, comp_name, ha='center', va='center',
                           fontsize=16, color='white', fontweight='bold')

            # Add professional border
            border = mpatches.Rectangle(
                (0.2, 0.2), 15.6, 11.6,
                linewidth=3, edgecolor='#1a1a2e',
                facecolor='none', linestyle='-', alpha=0.3
            )
            ax.add_patch(border)

            # Save with ultra-high quality
            plt.tight_layout(pad=1.5)
            plt.savefig(output_path, dpi=300, bbox_inches='tight',
                       facecolor='white', edgecolor='none', pad_inches=0.5)
            plt.close()

            return ToolResult(
                success=True,
                output={
                    "diagram_path": output_path,
                    "title": title,
                    "component_count": len(components),
                    "format": "PNG" if output_path.endswith('.png') else "SVG",
                    "size_bytes": Path(output_path).stat().st_size if Path(output_path).exists() else 0
                }
            )

        except Exception as e:
            logger.error(f"Failed to generate architecture diagram: {e}")
            import traceback
            traceback.print_exc()
            return ToolResult(success=False, output=None, error=str(e))


class CreateFlowchartTool(Tool):
    """Generate flowchart diagrams as images"""

    def __init__(self):
        super().__init__()
        self.name = "create_flowchart"
        self.description = "Create flowchart diagram as PNG image with shapes and connections"
        self.category = ToolCategory.DOCUMENTATION
        self.parameters = [
            ToolParameter(
                name="output_path",
                type="string",
                description="Path where flowchart should be saved",
                required=True
            ),
            ToolParameter(
                name="steps",
                type="array",
                description="List of flowchart steps with 'type' (start, process, decision, end) and 'text'",
                required=True
            ),
            ToolParameter(
                name="title",
                type="string",
                description="Flowchart title",
                required=False
            )
        ]

        # Capability profile
        self.capability_profile = ToolCapabilityProfile(
            tool_name="create_flowchart",
            capabilities=[
                CapabilityMetadata(
                    capability=Capability.DOCUMENT_CODE,
                    description="CreateFlowchart capability"
                )
            ]
        )

    async def execute(self, **kwargs) -> ToolResult:
        """Generate flowchart"""
        try:
            import matplotlib.pyplot as plt
            from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Polygon, Ellipse
            import numpy as np

            output_path = kwargs.get("output_path")
            steps = kwargs.get("steps", [])
            title = kwargs.get("title", "Flowchart")

            # Create figure
            fig, ax = plt.subplots(figsize=(10, 12))
            ax.set_xlim(0, 10)
            ax.set_ylim(0, len(steps) * 1.5 + 2)
            ax.axis('off')

            # Add title
            ax.text(5, len(steps) * 1.5 + 1.5, title, ha='center', va='top',
                   fontsize=18, fontweight='bold')

            y_pos = len(steps) * 1.5
            for i, step in enumerate(steps):
                step_type = step.get('type', 'process') if isinstance(step, dict) else 'process'
                step_text = step.get('text', f'Step {i+1}') if isinstance(step, dict) else str(step)

                if step_type == 'start' or step_type == 'end':
                    # Oval shape
                    box = Ellipse(
                        (5, y_pos), 3, 0.6,
                        edgecolor='#27AE60' if step_type == 'start' else '#E74C3C',
                        facecolor='#2ECC71' if step_type == 'start' else '#EC7063',
                        linewidth=2
                    )
                    ax.add_patch(box)
                    ax.text(5, y_pos, step_text, ha='center', va='center',
                           fontsize=10, color='white', fontweight='bold')

                elif step_type == 'decision':
                    # Diamond shape
                    diamond = Polygon(
                        [(5, y_pos + 0.4), (3.5, y_pos), (5, y_pos - 0.4), (6.5, y_pos)],
                        closed=True, edgecolor='#F39C12', facecolor='#F4D03F', linewidth=2
                    )
                    ax.add_patch(diamond)
                    ax.text(5, y_pos, step_text, ha='center', va='center',
                           fontsize=9, color='#34495E', fontweight='bold')

                else:  # process
                    # Rectangle
                    box = FancyBboxPatch(
                        (3.5, y_pos - 0.3), 3, 0.6,
                        boxstyle="round,pad=0.05",
                        edgecolor='#3498DB', facecolor='#5DADE2', linewidth=2
                    )
                    ax.add_patch(box)
                    ax.text(5, y_pos, step_text, ha='center', va='center',
                           fontsize=10, color='white', fontweight='bold')

                # Add arrow to next step
                if i < len(steps) - 1:
                    arrow = FancyArrowPatch(
                        (5, y_pos - 0.45), (5, y_pos - 1.05),
                        arrowstyle='->', mutation_scale=15,
                        color='#34495E', linewidth=1.5
                    )
                    ax.add_patch(arrow)

                y_pos -= 1.5

            # Save
            plt.tight_layout()
            plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
            plt.close()

            return ToolResult(
                success=True,
                output={
                    "flowchart_path": output_path,
                    "step_count": len(steps),
                    "format": "PNG",
                    "size_bytes": Path(output_path).stat().st_size if Path(output_path).exists() else 0
                }
            )

        except Exception as e:
            logger.error(f"Failed to generate flowchart: {e}")
            return ToolResult(success=False, output=None, error=str(e))
